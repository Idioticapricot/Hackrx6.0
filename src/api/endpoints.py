# src/api/endpoints.py
"""API endpoint implementations"""

import os
import asyncio
import httpx
import tiktoken
import base64
from concurrent.futures import ThreadPoolExecutor
from fastapi import HTTPException
from typing import List, Dict, Any
from pathlib import Path
from pptx import Presentation

from ..models.schemas import HackathonRequest, LegalAnalysisRequest
from ..ai import get_embed_model, get_reranker_model, STRICT_CONTEXT_SYSTEM_PROMPT, FULL_TEXT_SYSTEM_PROMPT, IMAGE_SYSTEM_PROMPT, get_document_system_prompt, get_legal_prompt
from ..ai.gemini_client import generate_text_with_gemini, analyze_image_with_gemini
from ..document_processing import load_and_process_document, get_processed_data
from ..utils import (
    check_for_secret_token_url, log_request_and_response, clean_text_for_llm
)
from ..utils.legal_analyzer import legal_analyzer
from ..ai.prompts import get_legal_prompt
from ..utils.terminal_ui import log_request_start, log_request_complete, create_request_progress
from ..core.config import GEMINI_API_KEY, OPENROUTER_API_KEY, YOUR_SITE_URL, YOUR_SITE_NAME, SMALL_DOC_TOKEN_THRESHOLD, USER_ID, API_BASE_URL, CACHE_DIR, config

# Initialize thread executor
executor = ThreadPoolExecutor(max_workers=os.cpu_count())

# Semaphores for rate limiting
llm_semaphore = asyncio.Semaphore(config.api.llm_semaphore_limit)
image_api_semaphore = asyncio.Semaphore(config.api.image_semaphore_limit)

async def rate_limited_llm_call(prompt: str, system_prompt: str = "") -> str:
    """Make rate-limited Gemini API call"""
    async with llm_semaphore:
        return await generate_text_with_gemini(prompt, system_prompt)

async def rate_limited_image_call(image_url: str, prompt: str, system_prompt: str = "") -> str:
    """Make rate-limited Gemini Vision API call"""
    async with image_api_semaphore:
        return await analyze_image_with_gemini(image_url, prompt, system_prompt)

def count_tokens(text: str) -> int:
    """Count tokens in text"""
    try:
        encoding = tiktoken.get_encoding(config.api.tiktoken_encoding)
        return len(encoding.encode(text))
    except Exception:
        return len(text.split())

def is_image_url(url: str) -> bool:
    """Check if URL points to an image file"""
    image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg', '.tiff', '.ico']
    url_lower = url.lower()
    
    # Extract the path part before query parameters
    url_path = url_lower.split('?')[0]
    
    return any(url_path.endswith(ext) for ext in image_extensions)

def is_pptx_url(url: str) -> bool:
    """Check if URL points to a PPTX file"""
    url_lower = url.lower()
    return url_lower.endswith(".pptx") or ".pptx?" in url_lower

def load_pptx_file(file_path: str) -> List[Dict[str, Any]]:
    """Extract text and images from PPTX file"""
    presentation_content = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slide_images = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f"slide_{i+1}_img_{len(slide_images)+1}.{image.ext}"
                    image_path = CACHE_DIR / image_filename
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    slide_images.append(str(image_path))

            presentation_content.append({
                "slide_number": i + 1,
                "text": "\n".join(slide_text),
                "image_paths": slide_images
            })
    except Exception as e:
        print(f"❌ Error loading PPTX: {e}")
    return presentation_content

async def analyze_local_image_with_gemini(image_path: str, question: str) -> str:
    """Analyze local image using Gemini Vision API"""
    try:
        # Upload image to temporary service to get URL (simplified approach)
        async with httpx.AsyncClient() as client:
            with open(image_path, "rb") as f:
                image_data = f.read()
            base64_image = "data:image/png;base64," + base64.b64encode(image_data).decode()
            
            # Upload to get public URL
            upload_payload = {"image_base64": base64_image, "filename": os.path.basename(image_path), "user_id": USER_ID}
            upload_response = await client.post(f"{API_BASE_URL}/upload-image", json=upload_payload, timeout=60)
            upload_response.raise_for_status()
            upload_json = upload_response.json()
            public_url = upload_json["url"]
            storage_path = upload_json["storage_path"]
            
            # Analyze with Gemini
            result = await analyze_image_with_gemini(public_url, question)
            
            # Clean up
            delete_payload = {"storage_path": storage_path, "user_id": USER_ID}
            await client.post(f"{API_BASE_URL}/delete-image", json=delete_payload, timeout=60)
            
            return result
            
    except Exception as e:
        return f"Error analyzing image: {str(e)}"

async def answer_questions_from_pptx(questions: List[str], pptx_content: List[Dict[str, Any]]) -> List[str]:
    """Builds comprehensive context from PPTX text and image analysis, then answers questions"""
    image_descriptions = {}
    image_analysis_tasks = []
    
    # Define the strict description prompt
    description_prompt = "Describe this image in extreme detail. Transcribe all text, numbers, and symbols verbatim, exactly as they appear. Do not interpret, summarize, infer meaning, or use any external knowledge. Your output should be a literal transcription of the visual content."
    
    # Collect unique images for analysis
    for slide in pptx_content:
        for image_path in slide["image_paths"]:
            if image_path not in image_descriptions:
                image_descriptions[image_path] = None  # Placeholder
                task = analyze_local_image_with_gemini(
                    image_path=image_path,
                    question=description_prompt
                )
                image_analysis_tasks.append((image_path, task))
    
    # Process all images in parallel
    if image_analysis_tasks:
        print(f"🖼️ Analyzing {len(image_analysis_tasks)} unique images from PPTX...")
        results = await asyncio.gather(*[task for _, task in image_analysis_tasks])
        
        for i, (path, _) in enumerate(image_analysis_tasks):
            image_descriptions[path] = results[i]
        print("✅ Image analysis complete.")
    
    # Build comprehensive context
    full_context_parts = []
    for slide in pptx_content:
        full_context_parts.append(f"--- Slide {slide['slide_number']} ---")
        if slide['text'].strip():
            full_context_parts.append("Text on slide:")
            full_context_parts.append(slide['text'])

        image_descs_on_slide = [image_descriptions[p] for p in slide["image_paths"] if p in image_descriptions]
        if any(image_descs_on_slide):
             full_context_parts.append("\nImage Descriptions on slide:")
             for desc in image_descs_on_slide:
                   full_context_parts.append(f"- {desc}")
        full_context_parts.append("\n")

    full_context = "\n".join(full_context_parts)
    
    # Answer questions using combined context
    return await answer_questions_from_context(questions, full_context, ".pptx")

async def answer_questions_from_context(questions: List[str], context: str, file_extension: str = None) -> List[str]:
    """Answer questions using provided context with legal-specific prompts"""
    # Use legal simplification prompt for better legal document understanding
    legal_prompt = get_legal_prompt('simplify')
    system_prompt = f"{legal_prompt}\n\nWhen answering questions, provide clear, practical explanations that help users understand their legal rights and obligations."
    
    print(f"💬 Processing {len(questions)} questions with context ({len(context)} chars)...")
    
    tasks = []
    for i, question in enumerate(questions):
        prompt = f"Document context:\n{context[:6000]}{'...' if len(context) > 6000 else ''}\n\nQuestion: {question}\n\nAnswer in simple terms:"
        print(f"🔍 [{i+1}/{len(questions)}] Question: {question[:50]}...")
        tasks.append(rate_limited_llm_call(prompt, system_prompt))
    
    return await asyncio.gather(*tasks)

async def process_image_question(image_url: str, question: str) -> str:
    """Process question about an image using Gemini Vision API"""
    return await rate_limited_image_call(image_url, question, IMAGE_SYSTEM_PROMPT)

async def retrieve_contexts_for_questions(questions: List[str], top_k: int = 10) -> List[str]:
    """Retrieve contexts for multiple questions using batch embeddings"""
    faiss_index, processed_texts, processed_metadatas = get_processed_data()
    
    if faiss_index is None or not processed_texts:
        print("⚠️ No processed data available for retrieval")
        return [""] * len(questions)

    print(f"🔍 Batch processing {len(questions)} questions...")
    
    embed_model = get_embed_model()
    reranker_model = get_reranker_model()
    
    # Batch embed all questions at once (CPU optimized for MacBook)
    print("🧠 Batch embedding all questions (MacBook CPU optimized)...")
    question_embeddings = embed_model.encode(questions, batch_size=min(config.models.batch_size, len(questions)))
    question_embeddings_np = question_embeddings.astype("float32")
    
    # Process all questions in parallel
    async def process_single_question_context(i, question, question_embedding):
        print(f"🔎 [{i+1}/{len(questions)}] Searching for: {question[:50]}...")
        
        # Search for similar chunks
        distances, indices = faiss_index.search(question_embedding.reshape(1, -1), min(top_k * 2, len(processed_texts)))
        
        candidates = []
        for j, idx in enumerate(indices[0]):
            if idx < len(processed_texts):
                candidates.append({
                    'text': processed_texts[idx],
                    'metadata': processed_metadatas[idx],
                    'distance': distances[0][j]
                })
        
        if not candidates:
            return ""
        
        # Rerank candidates with GPU memory management
        query_text_pairs = [(question, candidate['text']) for candidate in candidates]
        try:
            rerank_scores = reranker_model.predict(query_text_pairs)
        except RuntimeError as e:
            if "out of memory" in str(e).lower():
                print(f"⚠️ GPU memory error during reranking, using distance-based ranking...")
                # Fallback to distance-based ranking
                rerank_scores = [-candidate['distance'] for candidate in candidates]
            else:
                raise e
        
        for j, score in enumerate(rerank_scores):
            candidates[j]['rerank_score'] = score
        
        candidates.sort(key=lambda x: x['rerank_score'], reverse=True)
        top_candidates = candidates[:top_k]
        
        context_parts = []
        for candidate in top_candidates:
            clean_text = clean_text_for_llm(candidate['text'], candidate['metadata'])
            if clean_text:
                context_parts.append(clean_text)
        
        return "\n\n".join(context_parts)
    
    # Run all context retrievals in parallel
    import asyncio
    loop = asyncio.get_event_loop()
    tasks = []
    for i, (question, question_embedding) in enumerate(zip(questions, question_embeddings_np)):
        task = loop.run_in_executor(None, lambda i=i, q=question, emb=question_embedding: 
                                   asyncio.run(process_single_question_context(i, q, emb)))
        tasks.append(task)
    
    contexts = await asyncio.gather(*tasks)
    
    print(f"✅ Batch processing complete for {len(questions)} questions")
    return contexts

async def process_questions_with_rag(questions: List[str], doc_url: str) -> List[str]:
    """Process multiple questions using batch RAG pipeline with legal focus"""
    # Use legal-specific prompts for better legal document Q&A
    legal_prompt = get_legal_prompt('simplify')
    system_prompt = f"{legal_prompt}\n\nFocus on practical legal implications and user rights when answering questions."
    
    # Batch retrieve contexts for all questions
    contexts = await retrieve_contexts_for_questions(questions, top_k=10)
    
    # Process each question with its context
    tasks = []
    for question, context in zip(questions, contexts):
        print(f"🔍 Processing question: {question[:50]}... (context: {len(context)} chars)")
        if not context:
            tasks.append(asyncio.create_task(asyncio.sleep(0, result="I couldn't find relevant information in the document to answer this question.")))
        else:
            prompt = f"Document context:\n{context[:6000]}{'...' if len(context) > 6000 else ''}\n\nQuestion: {question}\n\nAnswer in simple terms:"
            tasks.append(rate_limited_llm_call(prompt, system_prompt))
    
    return await asyncio.gather(*tasks)

async def process_questions_with_full_text(questions: List[str], full_text: str, doc_url: str) -> List[str]:
    """Process multiple questions using full text with legal-specific prompts"""
    # Use legal-specific prompts for better legal document understanding
    legal_prompt = get_legal_prompt('simplify')
    system_prompt = f"{legal_prompt}\n\nProvide practical, actionable answers that help users understand their legal rights and obligations."
    
    print(f"💬 Processing {len(questions)} questions with full text ({len(full_text)} chars)...")
    
    tasks = []
    for i, question in enumerate(questions):
        prompt = f"Document:\n{full_text[:8000]}{'...' if len(full_text) > 8000 else ''}\n\nQuestion: {question}\n\nAnswer in simple terms:"
        print(f"🔍 [{i+1}/{len(questions)}] Question: {question[:50]}...")
        tasks.append(rate_limited_llm_call(prompt, system_prompt))
    
    return await asyncio.gather(*tasks)

async def hackathon_endpoint(request: HackathonRequest):
    """Main hackathon API endpoint with progress tracking"""
    import uuid
    import time
    
    doc_url = str(request.documents)
    questions = request.questions
    start_time = time.time()
    
    # Log request start
    log_request_start(doc_url, len(questions))
    
    # Create unique request ID and progress tracker
    request_id = str(uuid.uuid4())[:8]
    progress = create_request_progress(request_id, len(questions))
    
    if hasattr(progress, 'start'):
        progress.start()
    
    try:
        progress.update_step("🔍 Analyzing request", "Validating input and checking document type")
        progress.set_main_progress(10)
        
        # Check for secret token URLs first (in questions or document URL)
        has_secret_token, secret_token_result = check_for_secret_token_url(questions)
        
        # Also check if the document URL itself is a secret token URL
        if not has_secret_token and "get-secret-token" in doc_url:
            try:
                from ..utils.secret_token import fetch_secret_token_from_url
                secret_token_result = fetch_secret_token_from_url(doc_url)
                has_secret_token = True
            except Exception as e:
                secret_token_result = f"Error processing secret token URL: {str(e)}"
                has_secret_token = True
        
        if has_secret_token:
            progress.update_step("🔐 Secret Token", "Processing secret token")
            progress.set_main_progress(80)
            
            duration = time.time() - start_time
            progress.complete()
            log_request_complete(True, duration)
            
            # Return the secret token result for all questions
            answers = [secret_token_result] * len(questions)
            log_request_and_response({"documents": doc_url, "questions": questions}, answers)
            return {"answers": answers}
        
        progress.set_main_progress(20)

        # Check if the URL is a flight submission URL
        if "FinalRound4Submission" in doc_url:
            try:
                progress.update_step("✈️ Flight Processing", "Fetching flight information")
                progress.set_main_progress(70)
                from ..utils import flight_checker
                flight_api_url = config.api.flight_api_url
                flight_number = flight_checker.get_flight_number(flight_api_url)
                answers = [f"Your flight number is {flight_number}" for _ in questions]
                
                duration = time.time() - start_time
                progress.complete()
                log_request_complete(True, duration)
                log_request_and_response({"documents": doc_url, "questions": questions}, answers)
                return {"answers": answers}
            except Exception as e:
                duration = time.time() - start_time
                progress.fail(str(e))
                log_request_complete(False, duration)
                raise HTTPException(status_code=500, detail=f"Flight processing error: {str(e)}")

        # Check if the URL is a PPTX file
        elif is_pptx_url(doc_url):
            print("🗺️ PPTX file detected. Using slide analysis pipeline.")
            local_path = None
            pptx_content = []
            try:
                from ..document_processing.loaders import download_file
                local_path = download_file(doc_url)
                if not local_path:
                    raise HTTPException(status_code=500, detail="Could not download PPTX document.")
                
                pptx_content = load_pptx_file(local_path)
                if not pptx_content:
                    raise HTTPException(status_code=500, detail="Could not parse the PPTX document.")
                
                answers = await answer_questions_from_pptx(questions, pptx_content)
                log_request_and_response({"documents": doc_url, "questions": questions}, answers)
                return {"answers": answers}
            
            except Exception as e:
                if isinstance(e, HTTPException): 
                    raise e
                raise HTTPException(status_code=500, detail=f"PPTX processing error: {str(e)}")
            finally:
                if local_path and os.path.exists(local_path):
                    for slide in pptx_content:
                        for img_path in slide.get("image_paths", []):
                            if os.path.exists(img_path):
                                os.remove(img_path)
                    os.remove(local_path)
                    print(f"🗑️ Temporary PPTX file and extracted images cleaned up.")

        # Check if the URL is an image
        elif is_image_url(doc_url):
            progress.update_step("📥 Loading document", "Image URL detected")
            progress.set_main_progress(40)
            
            print(f"🖼️ Detected image URL: {doc_url}")
            print("📸 Processing image questions directly with vision API...")
            
            try:
                progress.update_step("❓ Processing", "Processing image questions")
                progress.set_main_progress(70)
                
                # Create tasks for parallel image processing
                async def process_single_image_question(i: int, question: str):
                    progress.update_question_progress(i, len(questions), f"Processing image: {question[:50]}...")
                    answer = await process_image_question(doc_url, question)
                    progress.update_question_progress(i + 1, len(questions), f"Completed image: {question[:50]}...")
                    return answer
                
                # Process all image questions concurrently
                tasks = [process_single_image_question(i, question) for i, question in enumerate(questions)]
                answers = await asyncio.gather(*tasks)
                
                progress.update_question_progress(len(questions), len(questions), "All questions completed")
                
                duration = time.time() - start_time
                progress.complete()
                log_request_complete(True, duration)
                log_request_and_response({"documents": doc_url, "questions": questions}, answers)
                return {"answers": answers}
                
            except Exception as e:
                duration = time.time() - start_time
                progress.fail(str(e))
                log_request_complete(False, duration)
                print(f"❌ Error processing image: {e}")
                raise HTTPException(status_code=500, detail=f"Image processing error: {str(e)}")

        try:
            # Load and process the document
            progress.update_step("📥 Loading document", "Downloading and parsing document")
            progress.set_main_progress(30)
            await asyncio.get_event_loop().run_in_executor(executor, load_and_process_document, doc_url)
            
            progress.update_step("📝 Processing document", "Extracting text and creating chunks")
            progress.set_main_progress(50)
            faiss_index, processed_texts, processed_metadatas = get_processed_data()
            
            # Check if document is small enough to bypass RAG
            full_text = "\n".join(processed_texts) if processed_texts else ""
            token_count = count_tokens(full_text)
            use_full_text = token_count <= SMALL_DOC_TOKEN_THRESHOLD
            
            if use_full_text:
                progress.update_step("🧠 Embedding content", "Skipped - using full text mode")
            else:
                progress.update_step("🧠 Embedding content", f"Embedded {len(processed_texts)} chunks")
            progress.set_main_progress(60)
            
            print(f"📊 Document stats: {token_count} tokens, {'using full text' if use_full_text else 'using RAG'}")

            progress.update_step("❓ Processing", "Generating answers for all questions")
            progress.set_main_progress(70)
            
            if use_full_text:
                # Batch process all questions with full text
                print("📄 Q&A Strategy: Full text (small document)")
                progress.update_question_progress(0, len(questions), "Processing all questions with full text...")
                answers = await process_questions_with_full_text(questions, full_text, doc_url)
            else:
                # Batch process all questions with RAG pipeline
                print("🔍 Q&A Strategy: RAG pipeline (targeted retrieval)")
                progress.update_question_progress(0, len(questions), "Processing all questions with RAG...")
                answers = await process_questions_with_rag(questions, doc_url)
                
                progress.update_question_progress(len(questions), len(questions), "All questions completed")
            
            progress.set_main_progress(90)
            
            duration = time.time() - start_time
            progress.complete()
            log_request_complete(True, duration)
            log_request_and_response({"documents": doc_url, "questions": questions}, answers)
            return {"answers": answers}

        except Exception as e:
            duration = time.time() - start_time
            progress.fail(str(e))
            log_request_complete(False, duration)
            print(f"❌ Error processing request: {e}")
            raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")
    
    except Exception as e:
        # Handle any outer exceptions
        if 'progress' in locals():
            duration = time.time() - start_time
            progress.fail(str(e))
            log_request_complete(False, duration)
        print(f"❌ Unexpected error: {e}")
        raise HTTPException(status_code=500, detail=f"Unexpected error: {str(e)}")

async def legal_analysis_endpoint(request: LegalAnalysisRequest):
    """Legal document analysis endpoint with complete analysis"""
    try:
        doc_url = str(request.document_url)
        language = request.language or 'en'
        
        # Load and process document
        await asyncio.get_event_loop().run_in_executor(executor, load_and_process_document, doc_url)
        faiss_index, processed_texts, processed_metadatas = get_processed_data()
        
        if not processed_texts:
            raise HTTPException(status_code=400, detail="Could not process document")
        
        # Combine all text for FULL CONTEXT analysis
        full_text = "\n\n".join(processed_texts)
        
        print(f"🔍 Using FULL CONTEXT strategy for legal analysis in {language}...")
        print(f"📄 Document size: {len(full_text)} chars, {len(processed_texts)} chunks")
        print("🧠 Sending ENTIRE document to Gemini for comprehensive analysis...")
        
        # Perform comprehensive legal analysis using FULL document context
        analysis = await legal_analyzer.analyze_document(full_text, language)
        
        # Add processing stats
        analysis['document_stats'] = {
            'total_length': len(full_text),
            'chunks_processed': len(processed_texts),
            'risk_count': len(analysis['risks']),
            'analysis_method': 'full_context_gemini'
        }
        
        print(f"✅ Full context analysis complete!")
        
        return {
            "analysis": analysis,
            "document_url": doc_url,
            "language": language,
            "processing_method": "full_context_gemini"
        }
        
    except Exception as e:
        print(f"❌ Legal analysis error: {e}")
        raise HTTPException(status_code=500, detail=f"Legal analysis error: {str(e)}")

async def simplify_document_endpoint(request: LegalAnalysisRequest):
    """Document simplification endpoint"""
    try:
        doc_url = str(request.document_url)
        language = request.language or 'en'
        
        # Load and process document
        await asyncio.get_event_loop().run_in_executor(executor, load_and_process_document, doc_url)
        faiss_index, processed_texts, processed_metadatas = get_processed_data()
        
        if not processed_texts:
            raise HTTPException(status_code=400, detail="Could not process document")
        
        # Combine all text
        full_text = "\n\n".join(processed_texts)
        
        # Simplify document
        simplified = await legal_analyzer.simplify_document(full_text, language)
        
        return {
            "original_text": full_text,
            "simplified_text": simplified,
            "document_url": doc_url,
            "language": language
        }
        
    except Exception as e:
        print(f"❌ Simplification error: {e}")
        raise HTTPException(status_code=500, detail=f"Simplification error: {str(e)}")

async def detect_risks_endpoint(request: LegalAnalysisRequest):
    """Risk detection endpoint"""
    try:
        doc_url = str(request.document_url)
        language = request.language or 'en'
        
        # Load and process document
        await asyncio.get_event_loop().run_in_executor(executor, load_and_process_document, doc_url)
        faiss_index, processed_texts, processed_metadatas = get_processed_data()
        
        if not processed_texts:
            raise HTTPException(status_code=400, detail="Could not process document")
        
        # Combine all text
        full_text = "\n\n".join(processed_texts)
        
        # Detect risks
        risks = await legal_analyzer.detect_risks(full_text, language)
        risk_summary = legal_analyzer.summarize_risks(risks)
        
        return {
            "risks": risks,
            "risk_summary": risk_summary,
            "document_url": doc_url,
            "language": language
        }
        
    except Exception as e:
        print(f"❌ Risk detection error: {e}")
        raise HTTPException(status_code=500, detail=f"Risk detection error: {str(e)}")